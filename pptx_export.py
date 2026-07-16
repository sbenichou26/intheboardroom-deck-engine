"""Editable PowerPoint export for In The Boardroom decks.

Two steps, kept separate on purpose:

1. deck_to_structure(): a cheap model call (sonnet, no web search) that reads the
   already-finished HTML deck and returns a JSON structure - titles, bullets,
   KPIs, chart data, tables, timeline entries, source lines. It only reformats
   content already present in the deck; it must not invent anything.

2. build_pptx(): turns that JSON into a native .pptx with python-pptx - editable
   text boxes, native tables and native charts (NOT flat images). The result can
   be opened and edited in PowerPoint/Keynote/Google Slides.

The HTML deck remains the source of truth; the PPTX is a convenience export.
"""
import json
import re
from io import BytesIO

# ---------------------------------------------------------------------------
# Step 1 - structure extraction (model call)
# ---------------------------------------------------------------------------

STRUCTURE_PROMPT = """You convert an already-written HTML investor deck into a JSON
structure for PowerPoint. Do NOT invent, add, round or change anything: use ONLY content
present in the HTML, and copy every figure verbatim.

Return ONLY a JSON object (no prose, no code fences) with this schema:
{
  "title": "the deck's main title",
  "subtitle": "the deck's subtitle or the asset name",
  "slides": [
    {
      "type": "section" | "content" | "kpi" | "chart" | "table" | "timeline" | "closing",
      "title": "slide title (short)",
      "subtitle": "optional one-line subclaim",
      "bullets": ["short bullet", "..."],
      "kpis": [{"value": "837", "label": "Revenue, EUR m, FY2024-25"}],
      "chart": {
        "kind": "column" | "bar" | "stacked" | "line",
        "unit": "EUR m",
        "categories": ["FY2022-23", "FY2024-25"],
        "series": [{"name": "Revenue", "values": [654, 837]}]
      },
      "table": {"headers": ["Col A", "Col B"], "rows": [["x", "y"], ["z", "w"]]},
      "timeline": [{"date": "2011", "text": "..."}],
      "source": "the slide's source-line text"
    }
  ]
}

Rules:
- Include only the keys relevant to each slide; omit the others.
- For a chart, read the plotted numbers from the inline SVG value labels. Only emit a
  "chart" when the slide genuinely has a numeric series; otherwise use bullets/kpi/table.
- Keep bullets short: split long paragraphs into separate bullets.
- Map the cover to a first "section" slide; map contact/closing to a "closing".
- Copy each figure exactly as written; never round or alter a number.
- Aim for one JSON slide per HTML slide, in the same order.
"""


def _slim_html(html: str) -> str:
    """Drop the <style> block and base64 image data before sending to the model:
    they carry no content and would waste a lot of tokens."""
    html = re.sub(r"<style.*?</style>", "", html, flags=re.S | re.I)
    html = re.sub(r"data:image/[^;]+;base64,[A-Za-z0-9+/=]+", "", html)
    return html


def _parse_structure(raw: str) -> dict:
    """Parse the model's JSON. If the output was truncated (a long deck can exceed
    the token budget), salvage it by keeping only the complete slide objects: try
    closing the JSON after each '}' from the end until it parses."""
    try:
        return json.loads(raw)
    except json.JSONDecodeError:
        i = len(raw)
        while True:
            i = raw.rfind("}", 0, i)
            if i == -1:
                break
            candidate = raw[:i + 1]
            for suffix in ("]}", "}]}", '"}]}'):
                try:
                    return json.loads(candidate + suffix)
                except json.JSONDecodeError:
                    continue
        raise


def deck_to_structure(client, deck_html: str, model: str = "claude-sonnet-5",
                      max_tokens: int = 16000) -> dict:
    """Ask the model to turn the finished HTML deck into the JSON structure above."""
    resp = client.messages.create(
        model=model,
        max_tokens=max_tokens,
        system=STRUCTURE_PROMPT,
        messages=[{"role": "user",
                   "content": "HTML deck to convert:\n\n" + _slim_html(deck_html)}],
    )
    raw = "".join(b.text for b in resp.content
                  if getattr(b, "type", None) == "text").strip()
    raw = re.sub(r"^```(json)?\s*", "", raw)
    raw = re.sub(r"```\s*$", "", raw)
    return _parse_structure(raw)


# ---------------------------------------------------------------------------
# Step 2 - build the .pptx
# ---------------------------------------------------------------------------

def build_pptx(structure: dict) -> BytesIO:
    """Render the JSON structure into a native, editable 16:9 .pptx."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

    NAVY = RGBColor(0x0D, 0x0A, 0x27)
    BLUE = RGBColor(0x12, 0x34, 0xFF)
    INK = RGBColor(0x1A, 0x1A, 0x2E)
    MUTE = RGBColor(0x6E, 0x6E, 0x8A)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    CARD = RGBColor(0xF6, 0xF7, 0xFB)

    KIND = {
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "stacked": XL_CHART_TYPE.COLUMN_STACKED,
        "bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "line": XL_CHART_TYPE.LINE_MARKERS,
    }

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def add_slide():
        return prs.slides.add_slide(blank)

    def fill_bg(slide, color):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = color

    def textbox(slide, left, top, width, height, text, size=18, color=INK,
                bold=False, italic=False, align=PP_ALIGN.LEFT):
        tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text or ""
        f = run.font
        f.size = Pt(size)
        f.bold = bold
        f.italic = italic
        f.color.rgb = color
        f.name = "Calibri"
        return tb

    def title_bar(slide, text):
        textbox(slide, 0.6, 0.4, 12.1, 0.9, text or "", size=26, color=NAVY, bold=True)
        rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      Inches(0.6), Inches(1.28), Inches(12.1), Pt(3))
        rule.fill.solid()
        rule.fill.fore_color.rgb = BLUE
        rule.line.fill.background()

    def source_line(slide, text):
        if text:
            textbox(slide, 0.6, 6.95, 12.1, 0.4, "Source: " + text,
                    size=9, color=MUTE, italic=True)

    def add_bullets(slide, bullets, top=1.7, size=16):
        tb = slide.shapes.add_textbox(Inches(0.7), Inches(top), Inches(11.9), Inches(4.9))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, b in enumerate(bullets or []):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            run = p.add_run()
            run.text = "•  " + str(b)
            run.font.size = Pt(size)
            run.font.color.rgb = INK
            run.font.name = "Calibri"
            p.space_after = Pt(8)

    def add_kpis(slide, kpis, top=1.9):
        kpis = (kpis or [])[:5]
        if not kpis:
            return
        gap = 0.3
        total_w = 12.1
        w = (total_w - gap * (len(kpis) - 1)) / len(kpis)
        x = 0.6
        for k in kpis:
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         Inches(x), Inches(top), Inches(w), Inches(2.0))
            box.fill.solid()
            box.fill.fore_color.rgb = CARD
            box.line.color.rgb = RGBColor(0xE4, 0xE4, 0xEF)
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = str(k.get("value", ""))
            r.font.size = Pt(30)
            r.font.bold = True
            r.font.color.rgb = BLUE
            r.font.name = "Calibri"
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            r2 = p2.add_run()
            r2.text = str(k.get("label", ""))
            r2.font.size = Pt(12)
            r2.font.color.rgb = MUTE
            r2.font.name = "Calibri"
            x += w + gap

    def add_chart(slide, chart, top=1.7):
        cats = chart.get("categories") or []
        series = chart.get("series") or []
        if not cats or not series:
            return
        cd = CategoryChartData()
        cd.categories = [str(c) for c in cats]
        added = 0
        for s in series:
            vals = []
            for v in s.get("values", []):
                try:
                    vals.append(float(str(v).replace(",", "").replace("%", "")))
                except (ValueError, TypeError):
                    vals.append(None)
            if any(v is not None for v in vals):
                cd.add_series(str(s.get("name", "Series")), vals)
                added += 1
        if added == 0:
            return
        ctype = KIND.get(chart.get("kind", "column"), XL_CHART_TYPE.COLUMN_CLUSTERED)
        gf = slide.shapes.add_chart(ctype, Inches(0.8), Inches(top),
                                    Inches(11.7), Inches(4.7), cd)
        ch = gf.chart
        ch.has_legend = added > 1
        if ch.has_legend:
            ch.legend.position = XL_LEGEND_POSITION.BOTTOM
            ch.legend.include_in_layout = False
        try:
            plot = ch.plots[0]
            plot.has_data_labels = True
            plot.data_labels.number_format = "0"
            plot.data_labels.number_format_is_linked = False
        except Exception:
            pass
        unit = chart.get("unit")
        if unit:
            textbox(slide, 0.8, top - 0.35, 6.0, 0.35, str(unit), size=11, color=MUTE)

    def add_table(slide, table, top=1.8):
        headers = table.get("headers") or []
        rows = table.get("rows") or []
        ncol = max([len(headers)] + [len(r) for r in rows] + [1])
        nrow = len(rows) + (1 if headers else 0)
        if nrow == 0:
            return
        height = min(5.0, 0.45 * nrow)
        gt = slide.shapes.add_table(nrow, ncol, Inches(0.6), Inches(top),
                                    Inches(12.1), Inches(height)).table
        ri = 0
        if headers:
            for ci in range(ncol):
                cell = gt.cell(0, ci)
                cell.text = str(headers[ci]) if ci < len(headers) else ""
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(11)
                        run.font.bold = True
            ri = 1
        for r in rows:
            for ci in range(ncol):
                cell = gt.cell(ri, ci)
                cell.text = str(r[ci]) if ci < len(r) else ""
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(11)
            ri += 1

    def add_timeline(slide, items, top=1.7):
        y = top
        for it in (items or [])[:10]:
            pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                          Inches(0.6), Inches(y), Inches(1.6), Inches(0.42))
            pill.fill.solid()
            pill.fill.fore_color.rgb = BLUE
            pill.line.fill.background()
            ptf = pill.text_frame
            pp = ptf.paragraphs[0]
            pp.alignment = PP_ALIGN.CENTER
            pr = pp.add_run()
            pr.text = str(it.get("date", ""))
            pr.font.size = Pt(12)
            pr.font.bold = True
            pr.font.color.rgb = WHITE
            pr.font.name = "Calibri"
            textbox(slide, 2.4, y, 10.2, 0.5, str(it.get("text", "")), size=13, color=INK)
            y += 0.52

    # --- title slide ---
    s = add_slide()
    fill_bg(s, NAVY)
    textbox(s, 0.9, 2.6, 11.5, 1.4, structure.get("title", "Discussion Materials"),
            size=40, color=WHITE, bold=True)
    if structure.get("subtitle"):
        textbox(s, 0.9, 4.0, 11.5, 0.8, structure.get("subtitle"), size=20, color=RGBColor(0xC7, 0xCC, 0xF5))
    textbox(s, 0.9, 6.9, 11.5, 0.4, "In The Boardroom · Discussion Materials",
            size=11, color=RGBColor(0xC7, 0xCC, 0xF5))

    # --- content slides ---
    for sl in structure.get("slides", []):
        stype = (sl.get("type") or "content").lower()
        try:
            if stype == "section":
                s = add_slide()
                fill_bg(s, NAVY)
                textbox(s, 0.9, 3.1, 11.5, 1.2, sl.get("title", ""),
                        size=32, color=WHITE, bold=True)
                if sl.get("subtitle"):
                    textbox(s, 0.9, 4.3, 11.5, 0.7, sl.get("subtitle"),
                            size=16, color=RGBColor(0xC7, 0xCC, 0xF5))
                continue

            if stype == "closing":
                s = add_slide()
                fill_bg(s, NAVY)
                textbox(s, 0.9, 3.0, 11.5, 1.2, sl.get("title", "Contact"),
                        size=30, color=WHITE, bold=True)
                for i, b in enumerate(sl.get("bullets", []) or []):
                    textbox(s, 0.9, 4.2 + i * 0.5, 11.5, 0.5, str(b), size=15,
                            color=RGBColor(0xC7, 0xCC, 0xF5))
                continue

            s = add_slide()
            fill_bg(s, WHITE)
            title_bar(s, sl.get("title", ""))
            if sl.get("subtitle"):
                textbox(s, 0.6, 1.35, 12.1, 0.5, sl.get("subtitle"),
                        size=14, color=MUTE, italic=True)
            body_top = 2.0 if sl.get("subtitle") else 1.7

            if stype == "kpi":
                add_kpis(s, sl.get("kpis"), top=body_top + 0.2)
                if sl.get("bullets"):
                    add_bullets(s, sl.get("bullets"), top=body_top + 2.5, size=14)
            elif stype == "chart":
                add_chart(s, sl.get("chart") or {}, top=body_top + 0.2)
                if sl.get("bullets"):
                    add_bullets(s, sl.get("bullets"), top=body_top, size=13)
            elif stype == "table":
                add_table(s, sl.get("table") or {}, top=body_top)
            elif stype == "timeline":
                add_timeline(s, sl.get("timeline"), top=body_top)
            else:  # content
                if sl.get("kpis"):
                    add_kpis(s, sl.get("kpis"), top=body_top)
                    add_bullets(s, sl.get("bullets"), top=body_top + 2.3, size=14)
                else:
                    add_bullets(s, sl.get("bullets"), top=body_top)

            source_line(s, sl.get("source"))
        except Exception:
            # A single malformed slide must never break the whole export.
            continue

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
